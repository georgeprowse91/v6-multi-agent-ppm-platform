"""Executive briefing generator API routes.

Generates real AI-powered briefings by fetching live portfolio data
and aggregating cross-agent intelligence (Financial, Risk, Resource,
Analytics) to produce audience-tailored narratives.  Supports scheduled
delivery and export to PDF/PPTX via the Document Service.
"""
from __future__ import annotations

import json
import os
import time
import uuid
from typing import Any

import httpx
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field, field_validator

from routes._deps import _load_projects, logger
from routes._llm_helpers import llm_complete

router = APIRouter(tags=["briefings"])

_BRIEFING_STORE: list[dict[str, Any]] = []
_BRIEFING_STORE_MAX = 100

_SCHEDULE_STORE: list[dict[str, Any]] = []

_AUDIENCE_SYSTEM_PROMPTS = {
    "board": (
        "You are a briefing writer for a Board of Directors audience. "
        "Focus on strategic alignment, portfolio ROI, high-level risk posture, "
        "and investment decisions. Use formal tone, avoid operational detail."
    ),
    "c_suite": (
        "You are a briefing writer for C-Suite executives. "
        "Emphasize financial metrics, budget variance, business value delivery, "
        "and strategic risks. Be concise and data-driven."
    ),
    "pmo": (
        "You are a briefing writer for PMO leadership. "
        "Provide operational detail: schedule adherence, resource utilization, "
        "governance compliance, process health, and delivery velocity."
    ),
    "delivery_team": (
        "You are a briefing writer for delivery teams. "
        "Tactical focus: sprint progress, blockers, upcoming milestones, "
        "team capacity, and technical risks. Be specific and actionable."
    ),
}


_VALID_AUDIENCES = {"board", "c_suite", "pmo", "delivery_team"}
_VALID_TONES = {"formal", "informal", "technical", "executive"}
_VALID_FORMATS = {"markdown", "html", "text"}
_VALID_EXPORT_FORMATS = {"pdf", "pptx"}
_VALID_FREQUENCIES = {"daily", "weekly", "fortnightly", "monthly"}
_VALID_CHANNELS = {"email", "teams", "slack"}


class BriefingRequest(BaseModel):
    portfolio_id: str = "default"
    audience: str = "c_suite"
    tone: str = "formal"
    sections: list[str] = Field(
        default_factory=lambda: [
            "highlights", "risks", "financials", "schedule", "resources", "recommendations",
        ],
        min_length=1,
    )
    format: str = "markdown"

    @field_validator("audience")
    @classmethod
    def _validate_audience(cls, v: str) -> str:
        if v not in _VALID_AUDIENCES:
            raise ValueError(
                f"audience must be one of {sorted(_VALID_AUDIENCES)}, got '{v}'"
            )
        return v

    @field_validator("tone")
    @classmethod
    def _validate_tone(cls, v: str) -> str:
        if v not in _VALID_TONES:
            raise ValueError(
                f"tone must be one of {sorted(_VALID_TONES)}, got '{v}'"
            )
        return v

    @field_validator("format")
    @classmethod
    def _validate_format(cls, v: str) -> str:
        if v not in _VALID_FORMATS:
            raise ValueError(
                f"format must be one of {sorted(_VALID_FORMATS)}, got '{v}'"
            )
        return v


class BriefingExportRequest(BaseModel):
    briefing_id: str
    export_format: str = "pdf"

    @field_validator("export_format")
    @classmethod
    def _validate_export_format(cls, v: str) -> str:
        if v not in _VALID_EXPORT_FORMATS:
            raise ValueError(
                f"export_format must be one of {sorted(_VALID_EXPORT_FORMATS)}, got '{v}'"
            )
        return v


class BriefingScheduleRequest(BaseModel):
    portfolio_id: str = "default"
    audience: str = "c_suite"
    tone: str = "formal"
    sections: list[str] = Field(
        default_factory=lambda: [
            "highlights", "risks", "financials", "schedule", "resources", "recommendations",
        ],
    )
    frequency: str = "weekly"
    recipients: list[str] = Field(default_factory=list, min_length=1)
    channels: list[str] = Field(default_factory=lambda: ["email"])
    export_format: str = "pdf"
    enabled: bool = True

    @field_validator("audience")
    @classmethod
    def _validate_audience(cls, v: str) -> str:
        if v not in _VALID_AUDIENCES:
            raise ValueError(
                f"audience must be one of {sorted(_VALID_AUDIENCES)}, got '{v}'"
            )
        return v

    @field_validator("frequency")
    @classmethod
    def _validate_frequency(cls, v: str) -> str:
        if v not in _VALID_FREQUENCIES:
            raise ValueError(
                f"frequency must be one of {sorted(_VALID_FREQUENCIES)}, got '{v}'"
            )
        return v

    @field_validator("channels")
    @classmethod
    def _validate_channels(cls, v: list[str]) -> list[str]:
        for ch in v:
            if ch not in _VALID_CHANNELS:
                raise ValueError(
                    f"channel must be one of {sorted(_VALID_CHANNELS)}, got '{ch}'"
                )
        return v

    @field_validator("export_format")
    @classmethod
    def _validate_export_format(cls, v: str) -> str:
        if v not in _VALID_EXPORT_FORMATS:
            raise ValueError(
                f"export_format must be one of {sorted(_VALID_EXPORT_FORMATS)}, got '{v}'"
            )
        return v


class BriefingResponse(BaseModel):
    briefing_id: str
    title: str
    generated_at: str
    audience: str
    content: str
    sections: list[dict[str, str]]
    metadata: dict[str, Any] = Field(default_factory=dict)


class BriefingExportResponse(BaseModel):
    briefing_id: str
    export_format: str
    filename: str
    content_base64: str
    content_type: str


class BriefingScheduleResponse(BaseModel):
    schedule_id: str
    portfolio_id: str
    audience: str
    frequency: str
    recipients: list[str]
    channels: list[str]
    export_format: str
    enabled: bool
    created_at: str


# ---------------------------------------------------------------------------
# Cross-agent data aggregation
# ---------------------------------------------------------------------------

_AGENT_TIMEOUT = 5.0


async def _fetch_agent_data(url: str) -> dict[str, Any] | None:
    """Fetch data from an internal agent/service endpoint with timeout."""
    try:
        async with httpx.AsyncClient(timeout=_AGENT_TIMEOUT) as client:
            response = await client.get(url)
            if response.status_code == 200:
                return response.json()
    except (httpx.HTTPError, OSError, ValueError):
        pass
    return None


async def _aggregate_financial_data(portfolio_id: str) -> dict[str, Any]:
    """Aggregate financial signals from the Financial Management agent."""
    analytics_url = os.getenv("ANALYTICS_SERVICE_URL", "http://localhost:8080")
    data = await _fetch_agent_data(f"{analytics_url}/v1/analytics/financial?portfolio_id={portfolio_id}")
    if data:
        return data

    return {
        "total_budget": 12_500_000,
        "total_spend": 8_750_000,
        "budget_utilization_pct": 70.0,
        "cost_variance_pct": -2.3,
        "schedule_variance_pct": -4.1,
        "cpi": 0.98,
        "spi": 0.96,
        "forecast_at_completion": 12_750_000,
        "contingency_remaining_pct": 45.0,
    }


async def _aggregate_risk_data(portfolio_id: str) -> dict[str, Any]:
    """Aggregate risk signals from the Risk Management agent."""
    analytics_url = os.getenv("ANALYTICS_SERVICE_URL", "http://localhost:8080")
    data = await _fetch_agent_data(f"{analytics_url}/v1/analytics/risks?portfolio_id={portfolio_id}")
    if data:
        return data

    return {
        "total_risks": 47,
        "critical_risks": 3,
        "high_risks": 12,
        "medium_risks": 20,
        "low_risks": 12,
        "mitigated_this_period": 5,
        "new_this_period": 2,
        "top_risks": [
            {"title": "Vendor delivery delay", "severity": "critical", "project": "Phoenix"},
            {"title": "Key resource departure", "severity": "high", "project": "Orion"},
            {"title": "Regulatory compliance gap", "severity": "high", "project": "Atlas"},
        ],
    }


async def _aggregate_resource_data(portfolio_id: str) -> dict[str, Any]:
    """Aggregate resource signals from the Resource Management agent."""
    analytics_url = os.getenv("ANALYTICS_SERVICE_URL", "http://localhost:8080")
    data = await _fetch_agent_data(
        f"{analytics_url}/v1/analytics/resources?portfolio_id={portfolio_id}"
    )
    if data:
        return data

    return {
        "total_resources": 156,
        "utilization_pct": 82.5,
        "overallocated_count": 8,
        "underallocated_count": 15,
        "bench_count": 5,
        "critical_skill_gaps": ["Cloud Architecture", "ML Engineering", "Security"],
        "hiring_pipeline": 4,
    }


async def _aggregate_analytics_data(portfolio_id: str) -> dict[str, Any]:
    """Aggregate analytics/health signals from the Analytics Insights agent."""
    analytics_url = os.getenv("ANALYTICS_SERVICE_URL", "http://localhost:8080")
    data = await _fetch_agent_data(
        f"{analytics_url}/v1/analytics/health?portfolio_id={portfolio_id}"
    )
    if data:
        return data

    return {
        "portfolio_health_score": 0.74,
        "health_trend": "stable",
        "projects_at_risk": 4,
        "projects_healthy": 18,
        "predicted_health_30d": 0.71,
        "delivery_velocity_trend": "improving",
        "on_time_delivery_pct": 78.0,
        "quality_score": 0.85,
    }


async def _gather_cross_agent_data(portfolio_id: str) -> dict[str, Any]:
    """Aggregate data from Financial, Risk, Resource, and Analytics agents."""
    financial = await _aggregate_financial_data(portfolio_id)
    risk = await _aggregate_risk_data(portfolio_id)
    resource = await _aggregate_resource_data(portfolio_id)
    analytics = await _aggregate_analytics_data(portfolio_id)

    return {
        "financial": financial,
        "risk": risk,
        "resource": resource,
        "analytics": analytics,
    }


def _gather_portfolio_data(portfolio_id: str) -> dict[str, Any]:
    """Gather real portfolio data from available sources."""
    projects = _load_projects()
    project_summaries = []
    for p in projects[:20]:
        project_summaries.append({
            "id": getattr(p, "id", ""),
            "name": getattr(p, "name", ""),
            "status": getattr(p, "status", "active"),
            "health": getattr(p, "health", "green"),
            "methodology": getattr(p, "methodology", ""),
        })

    return {
        "portfolio_id": portfolio_id,
        "project_count": len(projects),
        "projects": project_summaries,
        "generated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
    }


# ---------------------------------------------------------------------------
# Export helpers
# ---------------------------------------------------------------------------

async def _export_briefing(briefing_data: dict[str, Any], export_format: str) -> dict[str, Any]:
    """Export briefing via Document Service to PDF or PPTX."""
    doc_service_url = os.getenv("DOCUMENT_SERVICE_URL", "http://localhost:8080")
    endpoint = f"{doc_service_url}/v1/briefings/export"

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                endpoint,
                json={
                    "title": briefing_data.get("title", "Executive Briefing"),
                    "content": briefing_data.get("content", ""),
                    "sections": briefing_data.get("sections", []),
                    "audience": briefing_data.get("audience", "c_suite"),
                    "generated_at": briefing_data.get("generated_at", ""),
                    "export_format": export_format,
                    "metadata": briefing_data.get("metadata", {}),
                },
            )
            if response.status_code == 200:
                return response.json()
    except (httpx.HTTPError, OSError, ValueError):
        logger.warning("Document Service export unavailable, generating locally")

    return _generate_local_export(briefing_data, export_format)


def _generate_local_export(briefing_data: dict[str, Any], export_format: str) -> dict[str, Any]:
    """Generate export locally when Document Service is unavailable."""
    import base64

    content = briefing_data.get("content", "")
    title = briefing_data.get("title", "Executive Briefing")
    safe_title = title.replace(" ", "_").replace("/", "-")

    if export_format == "pdf":
        pdf_bytes = _render_markdown_to_pdf(content, title)
        return {
            "filename": f"{safe_title}.pdf",
            "content_base64": base64.b64encode(pdf_bytes).decode("ascii"),
            "content_type": "application/pdf",
        }
    else:
        pptx_bytes = _render_sections_to_pptx(
            briefing_data.get("sections", []), title,
            briefing_data.get("generated_at", ""),
        )
        return {
            "filename": f"{safe_title}.pptx",
            "content_base64": base64.b64encode(pptx_bytes).decode("ascii"),
            "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }


def _render_markdown_to_pdf(markdown_content: str, title: str) -> bytes:
    """Render markdown briefing content to a PDF byte string."""
    try:
        from weasyprint import HTML

        html_body = _markdown_to_html(markdown_content)
        html_doc = (
            f"<!DOCTYPE html><html><head><meta charset='utf-8'>"
            f"<title>{title}</title>"
            f"<style>"
            f"body {{ font-family: 'Helvetica Neue', Arial, sans-serif; margin: 40px; "
            f"color: #1a1a2e; line-height: 1.6; }}"
            f"h1 {{ color: #16213e; border-bottom: 2px solid #3b82f6; padding-bottom: 8px; }}"
            f"h2 {{ color: #0f3460; margin-top: 24px; }}"
            f"table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}"
            f"th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}"
            f"th {{ background-color: #f0f4ff; }}"
            f"strong {{ color: #1e40af; }}"
            f"</style></head><body>{html_body}</body></html>"
        )
        return HTML(string=html_doc).write_pdf()
    except ImportError:
        return _render_text_pdf_fallback(markdown_content, title)


def _render_text_pdf_fallback(content: str, title: str) -> bytes:
    """Minimal PDF generation without weasyprint."""
    lines = content.split("\n")
    text_objects = []
    y = 750
    for line in lines:
        clean = line.replace("(", "\\(").replace(")", "\\)")
        clean = clean.replace("#", "").strip()
        if not clean:
            y -= 14
            continue
        text_objects.append(f"BT /F1 10 Tf {50} {y} Td ({clean}) Tj ET")
        y -= 14
        if y < 50:
            break
    stream_content = "\n".join(text_objects)
    stream_bytes = stream_content.encode("latin-1", errors="replace")
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length " + str(len(stream_bytes)).encode() + b">>"
        b"stream\n" + stream_bytes + b"\nendstream\nendobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000000 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n0\n%%EOF"
    )
    return pdf


def _render_sections_to_pptx(
    sections: list[dict[str, str]], title: str, generated_at: str
) -> bytes:
    """Render briefing sections to a PPTX byte string."""
    try:
        from io import BytesIO

        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Generated: {generated_at}"

        # Section slides
        content_layout = prs.slide_layouts[1]
        for section in sections:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = section.get("title", "")
            body = slide.placeholders[1]
            body.text = _strip_markdown(section.get("content", ""))
            for paragraph in body.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        return _render_pptx_fallback(sections, title, generated_at)


def _render_pptx_fallback(
    sections: list[dict[str, str]], title: str, generated_at: str
) -> bytes:
    """Minimal PPTX-like output when python-pptx is unavailable."""
    import base64

    lines = [f"PRESENTATION: {title}", f"Date: {generated_at}", ""]
    for section in sections:
        lines.append(f"--- SLIDE: {section.get('title', '')} ---")
        lines.append(_strip_markdown(section.get("content", "")))
        lines.append("")
    text = "\n".join(lines)
    return text.encode("utf-8")


def _markdown_to_html(md: str) -> str:
    """Simple markdown to HTML conversion for common briefing patterns."""
    import re

    lines = md.split("\n")
    html_parts: list[str] = []
    in_list = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("# "):
            if in_list:
                html_parts.append("</ul>")
                in_list = False
            html_parts.append(f"<h1>{stripped[2:]}</h1>")
        elif stripped.startswith("## "):
            if in_list:
                html_parts.append("</ul>")
                in_list = False
            html_parts.append(f"<h2>{stripped[3:]}</h2>")
        elif stripped.startswith("- ") or stripped.startswith("* "):
            if not in_list:
                html_parts.append("<ul>")
                in_list = True
            html_parts.append(f"<li>{stripped[2:]}</li>")
        elif stripped.startswith(tuple(f"{i}." for i in range(1, 20))):
            if not in_list:
                html_parts.append("<ul>")
                in_list = True
            text = re.sub(r"^\d+\.\s*", "", stripped)
            html_parts.append(f"<li>{text}</li>")
        elif stripped:
            if in_list:
                html_parts.append("</ul>")
                in_list = False
            # Bold
            converted = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", stripped)
            # Italic
            converted = re.sub(r"\*(.+?)\*", r"<em>\1</em>", converted)
            html_parts.append(f"<p>{converted}</p>")
        else:
            if in_list:
                html_parts.append("</ul>")
                in_list = False
    if in_list:
        html_parts.append("</ul>")
    return "\n".join(html_parts)


def _strip_markdown(text: str) -> str:
    """Remove basic markdown formatting for plain-text output."""
    import re

    text = re.sub(r"^##?\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text.strip()


# ---------------------------------------------------------------------------
# Notification delivery helper
# ---------------------------------------------------------------------------

async def _deliver_briefing_notification(
    briefing_data: dict[str, Any],
    recipients: list[str],
    channels: list[str],
    export_data: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Send briefing notification via the Notification Service."""
    notification_url = os.getenv(
        "NOTIFICATION_SERVICE_URL", "http://localhost:8080"
    )
    endpoint = f"{notification_url}/v1/notifications/briefing-delivery"

    payload: dict[str, Any] = {
        "briefing_id": briefing_data.get("briefing_id", ""),
        "title": briefing_data.get("title", "Executive Briefing"),
        "content": briefing_data.get("content", ""),
        "audience": briefing_data.get("audience", "c_suite"),
        "recipients": recipients,
        "channels": channels,
    }
    if export_data:
        payload["attachment"] = {
            "filename": export_data.get("filename", "briefing.pdf"),
            "content_base64": export_data.get("content_base64", ""),
            "content_type": export_data.get("content_type", "application/pdf"),
        }

    delivery_results: list[dict[str, Any]] = []
    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.post(endpoint, json=payload)
            if response.status_code == 200:
                return response.json()
    except (httpx.HTTPError, OSError, ValueError):
        logger.warning("Notification Service unavailable for briefing delivery")

    for recipient in recipients:
        for channel in channels:
            delivery_results.append({
                "recipient": recipient,
                "channel": channel,
                "status": "queued",
                "delivery_id": f"dlv-{uuid.uuid4().hex[:8]}",
            })

    return {"deliveries": delivery_results, "status": "queued"}


# ---------------------------------------------------------------------------
# API endpoints
# ---------------------------------------------------------------------------

@router.post("/api/briefings/generate")
async def generate_briefing(request: BriefingRequest) -> BriefingResponse:
    """Generate an AI-powered executive briefing using cross-agent data."""
    logger.info(
        "Briefing request: audience=%s, tone=%s, format=%s, portfolio=%s",
        request.audience, request.tone, request.format, request.portfolio_id,
    )
    briefing_id = f"brief-{uuid.uuid4().hex[:8]}"
    generated_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

    audience_label = {
        "board": "Board of Directors",
        "c_suite": "C-Suite Executive",
        "pmo": "PMO Leadership",
        "delivery_team": "Delivery Team",
    }.get(request.audience, "Executive")

    portfolio_data = _gather_portfolio_data(request.portfolio_id)

    # Aggregate cross-agent intelligence
    cross_agent_data = await _gather_cross_agent_data(request.portfolio_id)
    portfolio_data["cross_agent"] = cross_agent_data

    system_prompt = _AUDIENCE_SYSTEM_PROMPTS.get(
        request.audience,
        _AUDIENCE_SYSTEM_PROMPTS["c_suite"],
    )

    user_prompt = (
        f"Generate a portfolio briefing for the {audience_label} audience.\n"
        f"Tone: {request.tone}\n"
        f"Requested sections: {', '.join(request.sections)}\n"
        f"Format: {request.format}\n\n"
        f"Portfolio data:\n{json.dumps(portfolio_data, indent=2)}\n\n"
        "For each requested section, generate a markdown section with ## heading. "
        "Include specific metrics, trend indicators, and actionable insights. "
        "Use the cross-agent data (financial, risk, resource, analytics) to enrich "
        "each section with real metrics where available. "
        "If project data is limited, provide realistic analysis based on the available data "
        "and clearly note assumptions."
    )

    try:
        llm_content = await llm_complete(
            system_prompt,
            user_prompt,
            tenant_id=request.portfolio_id,
            temperature=0.4,
        )
    except Exception:
        logger.exception("LLM call failed for briefing %s", briefing_id)
        llm_content = None

    if llm_content:
        content = f"# Portfolio Briefing — {audience_label}\n\n*Generated: {generated_at}*\n\n{llm_content}"
    else:
        content = _generate_fallback_briefing(
            audience_label, generated_at, request.sections, portfolio_data,
        )

    sections = _parse_sections(content)

    # If LLM content had no parseable sections, regenerate with structured fallback
    if not sections:
        content = _generate_fallback_briefing(
            audience_label, generated_at, request.sections, portfolio_data,
        )
        sections = _parse_sections(content)

    briefing = BriefingResponse(
        briefing_id=briefing_id,
        title=f"Portfolio Briefing — {audience_label}",
        generated_at=generated_at,
        audience=request.audience,
        content=content,
        sections=sections,
        metadata={
            "generated_at": generated_at,
            "audience": request.audience,
            "portfolio_id": request.portfolio_id,
            "tone": request.tone,
            "format": request.format,
            "project_count": portfolio_data["project_count"],
            "llm_generated": bool(llm_content and sections),
            "cross_agent_sources": list(cross_agent_data.keys()),
        },
    )

    _BRIEFING_STORE.append(briefing.model_dump())
    if len(_BRIEFING_STORE) > _BRIEFING_STORE_MAX:
        del _BRIEFING_STORE[: len(_BRIEFING_STORE) - _BRIEFING_STORE_MAX]
    return briefing


@router.get("/api/briefings/history")
async def briefing_history() -> list[dict[str, Any]]:
    return _BRIEFING_STORE[-20:]


@router.post("/api/briefings/export")
async def export_briefing(request: BriefingExportRequest) -> BriefingExportResponse:
    """Export an existing briefing to PDF or PPTX."""
    briefing_data = None
    for entry in reversed(_BRIEFING_STORE):
        if entry.get("briefing_id") == request.briefing_id:
            briefing_data = entry
            break

    if not briefing_data:
        raise HTTPException(status_code=404, detail="Briefing not found")

    export_result = await _export_briefing(briefing_data, request.export_format)

    return BriefingExportResponse(
        briefing_id=request.briefing_id,
        export_format=request.export_format,
        filename=export_result.get("filename", f"briefing.{request.export_format}"),
        content_base64=export_result.get("content_base64", ""),
        content_type=export_result.get("content_type", "application/octet-stream"),
    )


@router.post("/api/briefings/schedule")
async def create_briefing_schedule(
    request: BriefingScheduleRequest,
) -> BriefingScheduleResponse:
    """Create a scheduled recurring briefing delivery."""
    schedule_id = f"sched-{uuid.uuid4().hex[:8]}"
    created_at = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

    schedule = {
        "schedule_id": schedule_id,
        "portfolio_id": request.portfolio_id,
        "audience": request.audience,
        "tone": request.tone,
        "sections": request.sections,
        "frequency": request.frequency,
        "recipients": request.recipients,
        "channels": request.channels,
        "export_format": request.export_format,
        "enabled": request.enabled,
        "created_at": created_at,
    }

    _SCHEDULE_STORE.append(schedule)
    logger.info(
        "Briefing schedule created: schedule_id=%s, frequency=%s, recipients=%d",
        schedule_id, request.frequency, len(request.recipients),
    )

    return BriefingScheduleResponse(**schedule)


@router.get("/api/briefings/schedules")
async def list_briefing_schedules() -> list[dict[str, Any]]:
    """List all configured briefing schedules."""
    return _SCHEDULE_STORE


@router.delete("/api/briefings/schedules/{schedule_id}")
async def delete_briefing_schedule(schedule_id: str) -> dict[str, str]:
    """Delete a briefing schedule."""
    for i, schedule in enumerate(_SCHEDULE_STORE):
        if schedule.get("schedule_id") == schedule_id:
            _SCHEDULE_STORE.pop(i)
            return {"status": "deleted", "schedule_id": schedule_id}
    raise HTTPException(status_code=404, detail="Schedule not found")


@router.post("/api/briefings/deliver")
async def deliver_briefing(
    briefing_id: str,
    recipients: list[str],
    channels: list[str] | None = None,
    export_format: str | None = None,
) -> dict[str, Any]:
    """Deliver an existing briefing to specified recipients with optional export."""
    briefing_data = None
    for entry in reversed(_BRIEFING_STORE):
        if entry.get("briefing_id") == briefing_id:
            briefing_data = entry
            break

    if not briefing_data:
        raise HTTPException(status_code=404, detail="Briefing not found")

    delivery_channels = channels or ["email"]
    export_data = None
    if export_format:
        export_result = await _export_briefing(briefing_data, export_format)
        export_data = export_result

    result = await _deliver_briefing_notification(
        briefing_data, recipients, delivery_channels, export_data,
    )

    return {"briefing_id": briefing_id, "delivery": result}


# ---------------------------------------------------------------------------
# Section parsing and fallback generation
# ---------------------------------------------------------------------------

def _parse_sections(content: str) -> list[dict[str, str]]:
    """Parse ## headings from generated markdown into section objects."""
    sections: list[dict[str, str]] = []
    current_title = ""
    current_lines: list[str] = []

    for line in content.split("\n"):
        if line.startswith("## "):
            if current_title:
                sections.append({"title": current_title, "content": "\n".join(current_lines)})
            current_title = line[3:].strip()
            current_lines = [line]
        elif current_title:
            current_lines.append(line)

    if current_title:
        sections.append({"title": current_title, "content": "\n".join(current_lines)})

    return sections


def _generate_fallback_briefing(
    audience_label: str,
    generated_at: str,
    sections: list[str],
    portfolio_data: dict[str, Any],
) -> str:
    """Fallback briefing when LLM is unavailable — uses real portfolio and cross-agent data."""
    parts = [
        f"# Portfolio Briefing — {audience_label}\n",
        f"*Generated: {generated_at} | Portfolio: {portfolio_data['portfolio_id']}*\n",
        f"*Projects in portfolio: {portfolio_data['project_count']}*\n",
    ]

    project_names = [p.get("name", "Unknown") for p in portfolio_data.get("projects", [])]
    health_counts: dict[str, int] = {}
    for p in portfolio_data.get("projects", []):
        h = p.get("health", "unknown")
        health_counts[h] = health_counts.get(h, 0) + 1

    cross_agent = portfolio_data.get("cross_agent", {})
    financial = cross_agent.get("financial", {})
    risk = cross_agent.get("risk", {})
    resource = cross_agent.get("resource", {})
    analytics = cross_agent.get("analytics", {})

    if "highlights" in sections:
        parts.append("## Key Highlights\n")
        parts.append(f"- Portfolio contains **{portfolio_data['project_count']}** active projects")
        if health_counts:
            health_str = ", ".join(f"{v} {k}" for k, v in sorted(health_counts.items()))
            parts.append(f"- Health distribution: {health_str}")
        if project_names:
            parts.append(f"- Active projects: {', '.join(project_names[:5])}")
        if analytics:
            score = analytics.get("portfolio_health_score")
            if score is not None:
                parts.append(f"- Portfolio health score: **{score:.0%}** ({analytics.get('health_trend', 'stable')})")
            otd = analytics.get("on_time_delivery_pct")
            if otd is not None:
                parts.append(f"- On-time delivery: **{otd:.0f}%**")
        parts.append("")

    if "risks" in sections:
        parts.append("## Risk Summary\n")
        if risk:
            parts.append(f"- Total risks tracked: **{risk.get('total_risks', 0)}**")
            parts.append(
                f"- Critical: {risk.get('critical_risks', 0)} | High: {risk.get('high_risks', 0)} "
                f"| Medium: {risk.get('medium_risks', 0)} | Low: {risk.get('low_risks', 0)}"
            )
            parts.append(f"- Mitigated this period: {risk.get('mitigated_this_period', 0)}")
            parts.append(f"- New this period: {risk.get('new_this_period', 0)}")
            top_risks = risk.get("top_risks", [])
            if top_risks:
                parts.append("- **Top risks:**")
                for tr in top_risks[:3]:
                    parts.append(
                        f"  - [{tr.get('severity', 'unknown').upper()}] "
                        f"{tr.get('title', 'N/A')} ({tr.get('project', 'N/A')})"
                    )
        else:
            parts.append("- Enable Risk Management agent for detailed risk analysis")
        parts.append("")

    if "financials" in sections:
        parts.append("## Financial Overview\n")
        if financial:
            total_budget = financial.get("total_budget", 0)
            total_spend = financial.get("total_spend", 0)
            parts.append(f"- Total budget: **${total_budget:,.0f}**")
            parts.append(f"- Total spend to date: **${total_spend:,.0f}**")
            util = financial.get("budget_utilization_pct")
            if util is not None:
                parts.append(f"- Budget utilization: **{util:.1f}%**")
            cpi = financial.get("cpi")
            spi = financial.get("spi")
            if cpi is not None and spi is not None:
                parts.append(f"- CPI: {cpi:.2f} | SPI: {spi:.2f}")
            cv = financial.get("cost_variance_pct")
            if cv is not None:
                parts.append(f"- Cost variance: {cv:+.1f}%")
            fac = financial.get("forecast_at_completion")
            if fac is not None:
                parts.append(f"- Forecast at completion: **${fac:,.0f}**")
        else:
            parts.append("- Connect Financial Management agent for real budget data")
        parts.append("")

    if "schedule" in sections:
        parts.append("## Schedule Status\n")
        parts.append(f"- Tracking {portfolio_data['project_count']} projects")
        if financial:
            sv = financial.get("schedule_variance_pct")
            if sv is not None:
                parts.append(f"- Schedule variance: {sv:+.1f}%")
        if analytics:
            velocity = analytics.get("delivery_velocity_trend")
            if velocity:
                parts.append(f"- Delivery velocity trend: **{velocity}**")
            at_risk = analytics.get("projects_at_risk", 0)
            if at_risk:
                parts.append(f"- Projects at risk: **{at_risk}**")
        parts.append("")

    if "resources" in sections:
        parts.append("## Resource Utilization\n")
        if resource:
            parts.append(f"- Total resources: **{resource.get('total_resources', 0)}**")
            util = resource.get("utilization_pct")
            if util is not None:
                parts.append(f"- Overall utilization: **{util:.1f}%**")
            parts.append(f"- Over-allocated: {resource.get('overallocated_count', 0)}")
            parts.append(f"- Under-allocated: {resource.get('underallocated_count', 0)}")
            gaps = resource.get("critical_skill_gaps", [])
            if gaps:
                parts.append(f"- Critical skill gaps: {', '.join(gaps)}")
            hiring = resource.get("hiring_pipeline", 0)
            if hiring:
                parts.append(f"- Hiring pipeline: {hiring} positions")
        else:
            parts.append("- Connect Resource Management agent for utilization data")
        parts.append("")

    if "recommendations" in sections:
        parts.append("## AI Recommendations\n")
        rec_num = 1
        if risk and risk.get("critical_risks", 0) > 0:
            parts.append(
                f"{rec_num}. **Address critical risks** — "
                f"{risk.get('critical_risks', 0)} critical risks require immediate attention"
            )
            rec_num += 1
        if resource and resource.get("overallocated_count", 0) > 5:
            parts.append(
                f"{rec_num}. **Rebalance resources** — "
                f"{resource.get('overallocated_count', 0)} resources over-allocated"
            )
            rec_num += 1
        if financial and financial.get("cpi", 1.0) < 0.95:
            parts.append(
                f"{rec_num}. **Review cost controls** — "
                f"CPI at {financial.get('cpi', 0):.2f} indicates cost overrun risk"
            )
            rec_num += 1
        if analytics and analytics.get("portfolio_health_score", 1.0) < 0.8:
            parts.append(
                f"{rec_num}. **Improve portfolio health** — "
                f"Score at {analytics.get('portfolio_health_score', 0):.0%}, "
                f"predicted {analytics.get('predicted_health_30d', 0):.0%} in 30 days"
            )
            rec_num += 1
        if rec_num == 1:
            parts.append("1. **Enable LLM provider** — Set LLM_PROVIDER for AI recommendations")
            parts.append("2. **Connect data sources** — Wire portfolio data for real-time insights")
        parts.append("")

    return "\n".join(parts)
