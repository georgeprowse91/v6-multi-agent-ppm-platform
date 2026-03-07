"""Briefing document renderer — PDF and PPTX export for executive briefings.

Provides rendering capabilities for the Executive Briefing Generator,
converting markdown briefing content into professionally formatted
PDF and PowerPoint documents.
"""
from __future__ import annotations

import base64
import logging
import re
from typing import Any

logger = logging.getLogger("document-service.briefing-renderer")


def render_briefing(
    title: str,
    content: str,
    sections: list[dict[str, str]],
    audience: str,
    generated_at: str,
    export_format: str,
    metadata: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Render a briefing to the requested export format.

    Returns a dict with filename, content_base64, and content_type.
    """
    safe_title = re.sub(r"[^a-zA-Z0-9_-]", "_", title)

    if export_format == "pdf":
        pdf_bytes = render_briefing_pdf(title, content, audience, generated_at, metadata)
        return {
            "filename": f"{safe_title}.pdf",
            "content_base64": base64.b64encode(pdf_bytes).decode("ascii"),
            "content_type": "application/pdf",
        }
    elif export_format == "pptx":
        pptx_bytes = render_briefing_pptx(title, sections, audience, generated_at, metadata)
        return {
            "filename": f"{safe_title}.pptx",
            "content_base64": base64.b64encode(pptx_bytes).decode("ascii"),
            "content_type": (
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
        }
    else:
        raise ValueError(f"Unsupported export format: {export_format}")


# ---------------------------------------------------------------------------
# PDF rendering
# ---------------------------------------------------------------------------

def render_briefing_pdf(
    title: str,
    content: str,
    audience: str,
    generated_at: str,
    metadata: dict[str, Any] | None = None,
) -> bytes:
    """Render markdown briefing content to a branded PDF."""
    try:
        return _render_pdf_weasyprint(title, content, audience, generated_at, metadata)
    except ImportError:
        logger.info("weasyprint not available, using fallback PDF renderer")
        return _render_pdf_fallback(title, content, audience, generated_at)


def _render_pdf_weasyprint(
    title: str,
    content: str,
    audience: str,
    generated_at: str,
    metadata: dict[str, Any] | None = None,
) -> bytes:
    """Render PDF using WeasyPrint with professional styling."""
    from weasyprint import HTML

    html_body = _markdown_to_html(content)
    audience_color = {
        "board": "#1e3a5f",
        "c_suite": "#16213e",
        "pmo": "#0f3460",
        "delivery_team": "#1a365d",
    }.get(audience, "#16213e")

    project_count = (metadata or {}).get("project_count", "")
    meta_line = f"Generated: {generated_at}"
    if project_count:
        meta_line += f" | Projects: {project_count}"

    html_doc = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{title}</title>
<style>
@page {{
    size: A4;
    margin: 2.5cm;
    @top-right {{ content: "Confidential"; font-size: 9pt; color: #999; }}
    @bottom-center {{ content: counter(page) " / " counter(pages); font-size: 9pt; color: #999; }}
}}
body {{
    font-family: 'Helvetica Neue', Arial, sans-serif;
    color: #1a1a2e;
    line-height: 1.7;
    font-size: 11pt;
}}
.cover-title {{
    font-size: 28pt;
    font-weight: 700;
    color: {audience_color};
    margin-bottom: 8px;
    border-bottom: 3px solid #3b82f6;
    padding-bottom: 12px;
}}
.cover-meta {{
    font-size: 10pt;
    color: #666;
    margin-bottom: 32px;
}}
h1 {{
    font-size: 22pt;
    color: {audience_color};
    border-bottom: 2px solid #3b82f6;
    padding-bottom: 6px;
    margin-top: 28px;
}}
h2 {{
    font-size: 16pt;
    color: #0f3460;
    margin-top: 24px;
    border-left: 4px solid #3b82f6;
    padding-left: 12px;
}}
table {{
    border-collapse: collapse;
    width: 100%;
    margin: 12px 0;
}}
th, td {{
    border: 1px solid #ddd;
    padding: 8px 10px;
    text-align: left;
    font-size: 10pt;
}}
th {{
    background-color: #f0f4ff;
    font-weight: 600;
}}
ul, ol {{
    padding-left: 20px;
}}
li {{
    margin-bottom: 4px;
}}
strong {{
    color: #1e40af;
}}
.metric-highlight {{
    background: #f0f4ff;
    padding: 2px 6px;
    border-radius: 3px;
    font-weight: 600;
}}
</style>
</head>
<body>
<div class="cover-title">{title}</div>
<div class="cover-meta">{meta_line}</div>
{html_body}
</body>
</html>"""

    return HTML(string=html_doc).write_pdf()


def _render_pdf_fallback(
    title: str,
    content: str,
    audience: str,
    generated_at: str,
) -> bytes:
    """Minimal PDF generation when WeasyPrint is not available."""
    lines = content.split("\n")
    text_objects = []
    y = 750

    # Title
    clean_title = title.replace("(", "\\(").replace(")", "\\)")
    text_objects.append(f"BT /F1 16 Tf {50} {y} Td ({clean_title}) Tj ET")
    y -= 24

    # Date
    clean_date = generated_at.replace("(", "\\(").replace(")", "\\)")
    text_objects.append(f"BT /F1 9 Tf {50} {y} Td (Generated: {clean_date}) Tj ET")
    y -= 20

    for line in lines:
        if y < 50:
            break
        clean = line.replace("(", "\\(").replace(")", "\\)")
        clean = clean.replace("#", "").strip()
        if not clean:
            y -= 10
            continue
        font_size = 10
        if line.startswith("## "):
            font_size = 13
            y -= 6
        elif line.startswith("# "):
            font_size = 16
            y -= 8
        text_objects.append(
            f"BT /F1 {font_size} Tf {50} {y} Td ({clean}) Tj ET"
        )
        y -= 14

    stream_content = "\n".join(text_objects)
    stream_bytes = stream_content.encode("latin-1", errors="replace")
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length " + str(len(stream_bytes)).encode() + b">>"
        b"stream\n" + stream_bytes + b"\nendstream\nendobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000000 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n0\n%%EOF"
    )
    return pdf


# ---------------------------------------------------------------------------
# PPTX rendering
# ---------------------------------------------------------------------------

def render_briefing_pptx(
    title: str,
    sections: list[dict[str, str]],
    audience: str,
    generated_at: str,
    metadata: dict[str, Any] | None = None,
) -> bytes:
    """Render briefing sections to a branded PowerPoint presentation."""
    try:
        return _render_pptx_python_pptx(title, sections, audience, generated_at, metadata)
    except ImportError:
        logger.info("python-pptx not available, using fallback PPTX renderer")
        return _render_pptx_fallback(title, sections, generated_at)


def _render_pptx_python_pptx(
    title: str,
    sections: list[dict[str, str]],
    audience: str,
    generated_at: str,
    metadata: dict[str, Any] | None = None,
) -> bytes:
    """Render PPTX using python-pptx with professional layout."""
    from io import BytesIO

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    audience_colors = {
        "board": RGBColor(0x1E, 0x3A, 0x5F),
        "c_suite": RGBColor(0x16, 0x21, 0x3E),
        "pmo": RGBColor(0x0F, 0x34, 0x60),
        "delivery_team": RGBColor(0x1A, 0x36, 0x5D),
    }
    title_color = audience_colors.get(audience, RGBColor(0x16, 0x21, 0x3E))

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.color.rgb = title_color
        run.font.bold = True
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generated: {generated_at}"
    for run in subtitle.text_frame.paragraphs[0].runs:
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Content slides
    content_layout = prs.slide_layouts[1]
    for section in sections:
        slide = prs.slides.add_slide(content_layout)
        section_title = slide.shapes.title
        section_title.text = section.get("title", "")
        for run in section_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.color.rgb = title_color
            run.font.bold = True

        body = slide.placeholders[1]
        body.text = _strip_markdown(section.get("content", ""))
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_pptx_fallback(
    title: str,
    sections: list[dict[str, str]],
    generated_at: str,
) -> bytes:
    """Plain-text PPTX-like output when python-pptx is unavailable."""
    lines = [f"PRESENTATION: {title}", f"Date: {generated_at}", ""]
    for section in sections:
        lines.append(f"--- SLIDE: {section.get('title', '')} ---")
        lines.append(_strip_markdown(section.get("content", "")))
        lines.append("")
    return "\n".join(lines).encode("utf-8")


# ---------------------------------------------------------------------------
# Markdown helpers
# ---------------------------------------------------------------------------

def _markdown_to_html(md: str) -> str:
    """Convert basic markdown to HTML for PDF rendering."""
    lines = md.split("\n")
    html_parts: list[str] = []
    in_list = False
    in_ol = False

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("# "):
            _close_list(html_parts, in_list, in_ol)
            in_list = in_ol = False
            html_parts.append(f"<h1>{_inline_md(stripped[2:])}</h1>")
        elif stripped.startswith("## "):
            _close_list(html_parts, in_list, in_ol)
            in_list = in_ol = False
            html_parts.append(f"<h2>{_inline_md(stripped[3:])}</h2>")
        elif stripped.startswith("- ") or stripped.startswith("* "):
            if not in_list:
                html_parts.append("<ul>")
                in_list = True
            html_parts.append(f"<li>{_inline_md(stripped[2:])}</li>")
        elif re.match(r"^\d+\.\s", stripped):
            if not in_ol:
                _close_list(html_parts, in_list, False)
                in_list = False
                html_parts.append("<ol>")
                in_ol = True
            text = re.sub(r"^\d+\.\s*", "", stripped)
            html_parts.append(f"<li>{_inline_md(text)}</li>")
        elif stripped:
            _close_list(html_parts, in_list, in_ol)
            in_list = in_ol = False
            html_parts.append(f"<p>{_inline_md(stripped)}</p>")
        else:
            _close_list(html_parts, in_list, in_ol)
            in_list = in_ol = False

    _close_list(html_parts, in_list, in_ol)
    return "\n".join(html_parts)


def _close_list(parts: list[str], in_ul: bool, in_ol: bool) -> None:
    if in_ul:
        parts.append("</ul>")
    if in_ol:
        parts.append("</ol>")


def _inline_md(text: str) -> str:
    """Convert inline markdown (bold, italic) to HTML."""
    text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"\*(.+?)\*", r"<em>\1</em>", text)
    return text


def _strip_markdown(text: str) -> str:
    """Remove markdown formatting for plain-text output."""
    text = re.sub(r"^##?\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text.strip()
